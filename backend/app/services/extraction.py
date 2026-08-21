"""
app/services/extraction.py
──────────────────────────
PDF text extraction service using pdfplumber (with PyMuPDF fallback).

WHY pdfplumber?
  - Preserves visual layout, multi-column reading order, and table boundaries.
  - Better than naive text extractors which can scramble adjacent columns.

Garbage In, Garbage Out:
  - If a PDF is scanned/image-only or corrupt, this service detects it early
    and raises a descriptive ValueError rather than passing empty strings to RAG.
"""

import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path: str) -> str:
    """
    Extracts plain text from a PDF file located at `file_path`.

    Args:
        file_path: Absolute or relative path to the PDF on disk.

    Returns:
        A string containing all extracted text with page separators.

    Raises:
        ValueError: If file does not exist, is corrupt/password-protected,
                    or contains no extractable text (e.g. scanned image).
    """
    if not os.path.exists(file_path):
        raise ValueError(f"File not found on disk: {file_path}")

    # Primary: pdfplumber for layout-aware extraction
    extracted_pages = []
    try:
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            if not pdf.pages:
                raise ValueError("PDF file contains 0 pages.")

            for page_idx, page in enumerate(pdf.pages, start=1):
                # extract_text with layout preservation
                page_text = page.extract_text(layout=True)
                if page_text and page_text.strip():
                    extracted_pages.append(f"--- Page {page_idx} ---\n{page_text.strip()}")

    except ImportError:
        logger.warning("pdfplumber not installed, attempting fallback with PyMuPDF (fitz)")
        try:
            import fitz  # PyMuPDF fallback
            doc = fitz.open(file_path)
            if len(doc) == 0:
                raise ValueError("PDF file contains 0 pages.")
            for page_idx in range(len(doc)):
                page = doc.load_page(page_idx)
                text = page.get_text("text")
                if text and text.strip():
                    extracted_pages.append(f"--- Page {page_idx + 1} ---\n{text.strip()}")
            doc.close()
        except Exception as fitz_err:
            raise ValueError(f"Failed to extract PDF text with PyMuPDF: {str(fitz_err)}")
    except Exception as e:
        if isinstance(e, ValueError):
            raise e
        # Catch PDFSyntaxError, PasswordRequired, corrupt streams, etc.
        raise ValueError(f"Corrupt or unreadable PDF: {str(e)}")

    full_text = "\n\n".join(extracted_pages).strip()

    if not full_text:
        raise ValueError(
            "PDF contains no extractable text. The document may be scanned images or password-protected."
        )

    return full_text


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extracts plain text from a PPT/PPTX file located at `file_path`.
    Pulls text from all shapes/slides (titles, body text) and speaker notes.

    Args:
        file_path: Absolute or relative path to the PPTX on disk.

    Returns:
        A string containing all extracted text with slide separators.

    Raises:
        ValueError: If file does not exist, is corrupt/unreadable,
                    or contains no extractable text.
    """
    if not os.path.exists(file_path):
        raise ValueError(f"File not found on disk: {file_path}")

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Corrupt or unreadable PPT/PPTX file: {str(e)}")

    extracted_slides = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_elements = []

        # Extract text from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame_content = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_frame_content.append(para_text)
                if text_frame_content:
                    slide_text_elements.append("\n".join(text_frame_content))

        # Extract speaker notes if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_text_elements.append(f"[Speaker Notes]\n{notes_text}")

        slide_body = "\n".join(slide_text_elements).strip()
        if slide_body:
            extracted_slides.append(f"--- Slide {slide_idx} ---\n{slide_body}")

    full_text = "\n\n".join(extracted_slides).strip()

    if not full_text:
        raise ValueError(
            "Presentation contains no extractable text (e.g. it may consist entirely of image-only slides)."
        )

    return full_text


def extract_text(file_path: str, file_type: str) -> str:
    """
    Extracts text from the document at `file_path` based on its `file_type`.
    """
    ext = file_type.lower().strip().replace(".", "")
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext in ("pptx", "ppt"):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

